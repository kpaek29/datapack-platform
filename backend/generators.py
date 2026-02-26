"""
Output Generators Module
Creates PPT and Excel outputs from processed data using master template
"""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.chart import XL_CHART_TYPE
from pptx.chart.data import CategoryChartData
import pandas as pd
from openpyxl import Workbook
from openpyxl.styles import Font, PatternFill, Alignment, Border, Side
from openpyxl.chart import BarChart, LineChart, Reference
from pathlib import Path
from typing import Dict, List, Any
from datetime import datetime

# Path to master slide template
TEMPLATE_PATH = Path(__file__).parent.parent / "templates" / "master_template.pptx"

class PPTGenerator:
    """Generate PowerPoint data packs using master template"""
    
    def __init__(self, output_path: Path, template_path: Path = None):
        self.output_path = output_path
        
        # Load from master template if available
        template = template_path or TEMPLATE_PATH
        if template.exists():
            self.prs = Presentation(str(template))
            # Remove placeholder slides from template
            while len(self.prs.slides) > 0:
                rId = self.prs.slides._sldIdLst[0].rId
                self.prs.part.drop_rel(rId)
                del self.prs.slides._sldIdLst[0]
            self._using_template = True
        else:
            self.prs = Presentation()
            self.prs.slide_width = Inches(13.333)  # Widescreen
            self.prs.slide_height = Inches(7.5)
            self._using_template = False
        
        self._setup_layouts()
    
    def _setup_layouts(self):
        """Setup slide layout references from template"""
        layouts = self.prs.slide_layouts
        
        # Map layouts by name or index
        self.layout_title = None
        self.layout_content = None
        self.layout_section = None
        self.layout_two_content = None
        self.layout_title_only = None
        self.layout_blank = None
        
        for i, layout in enumerate(layouts):
            name = layout.name.lower()
            if 'title slide' in name or (i == 0 and 'title' in name):
                self.layout_title = layout
            elif 'title and content' in name:
                self.layout_content = layout
            elif 'section' in name:
                self.layout_section = layout
            elif 'two content' in name:
                self.layout_two_content = layout
            elif 'title only' in name:
                self.layout_title_only = layout
            elif 'blank' in name:
                self.layout_blank = layout
        
        # Fallback assignments
        if self._using_template and len(layouts) >= 5:
            self.layout_title = self.layout_title or layouts[0]
            self.layout_content = self.layout_content or layouts[1]
            self.layout_section = self.layout_section or layouts[2]
            self.layout_two_content = self.layout_two_content or layouts[3]
            self.layout_title_only = self.layout_title_only or layouts[4]
        
        # Ultimate fallback for non-template
        if len(layouts) >= 7:
            self.layout_blank = self.layout_blank or layouts[6]
        elif len(layouts) > 0:
            self.layout_blank = layouts[0]
    
    def _set_placeholder_text(self, slide, ph_type: str, text: str):
        """Set text in a placeholder by type"""
        from pptx.enum.shapes import PP_PLACEHOLDER
        
        ph_map = {
            'title': PP_PLACEHOLDER.TITLE,
            'body': PP_PLACEHOLDER.BODY,
            'subtitle': PP_PLACEHOLDER.SUBTITLE,
            'center_title': PP_PLACEHOLDER.CENTER_TITLE,
        }
        
        target_type = ph_map.get(ph_type)
        
        for shape in slide.placeholders:
            if target_type and shape.placeholder_format.type == target_type:
                shape.text = text
                return shape
            elif ph_type == 'title' and 'title' in shape.name.lower():
                shape.text = text
                return shape
            elif ph_type == 'body' and ('content' in shape.name.lower() or 'body' in shape.name.lower()):
                shape.text = text
                return shape
        return None
    
    def add_title_slide(self, title: str, subtitle: str = None):
        """Add a title slide using template layout"""
        layout = self.layout_title or self.layout_blank or self.prs.slide_layouts[0]
        slide = self.prs.slides.add_slide(layout)
        
        if self._using_template:
            self._set_placeholder_text(slide, 'title', title)
            if subtitle:
                self._set_placeholder_text(slide, 'subtitle', subtitle)
        else:
            # Fallback: manual text boxes
            title_box = slide.shapes.add_textbox(Inches(0.5), Inches(2.5), Inches(12.333), Inches(1))
            title_frame = title_box.text_frame
            title_para = title_frame.paragraphs[0]
            title_para.text = title
            title_para.font.size = Pt(44)
            title_para.font.bold = True
            
            if subtitle:
                sub_box = slide.shapes.add_textbox(Inches(0.5), Inches(3.7), Inches(12.333), Inches(0.5))
                sub_frame = sub_box.text_frame
                sub_para = sub_frame.paragraphs[0]
                sub_para.text = subtitle
                sub_para.font.size = Pt(20)
                sub_para.font.color.rgb = RGBColor(100, 100, 100)
        
        return slide
    
    def add_section_slide(self, title: str):
        """Add a section divider slide using template layout"""
        layout = self.layout_section or self.layout_blank or self.prs.slide_layouts[0]
        slide = self.prs.slides.add_slide(layout)
        
        if self._using_template:
            self._set_placeholder_text(slide, 'title', title)
        else:
            title_box = slide.shapes.add_textbox(Inches(0.5), Inches(3), Inches(12.333), Inches(1))
            title_frame = title_box.text_frame
            title_para = title_frame.paragraphs[0]
            title_para.text = title
            title_para.font.size = Pt(36)
            title_para.font.bold = True
        
        return slide
    
    def add_table_slide(self, title: str, df: pd.DataFrame, max_rows: int = 15):
        """Add a slide with a data table using template layout"""
        layout = self.layout_title_only or self.layout_blank or self.prs.slide_layouts[0]
        slide = self.prs.slides.add_slide(layout)
        
        # Set title via placeholder or fallback
        if self._using_template:
            self._set_placeholder_text(slide, 'title', title)
        else:
            title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12.333), Inches(0.6))
            title_frame = title_box.text_frame
            title_para = title_frame.paragraphs[0]
            title_para.text = title
            title_para.font.size = Pt(24)
            title_para.font.bold = True
        
        # Table - position below title
        df_display = df.head(max_rows)
        rows, cols = df_display.shape
        rows += 1  # Header
        
        # Use template-appropriate positioning
        table_width = Inches(9.2) if self._using_template else Inches(12.333)
        table = slide.shapes.add_table(rows, cols, Inches(0.4), Inches(1.2), table_width, Inches(5.5)).table
        
        # Header with template styling
        header_color = RGBColor(0x08, 0x46, 0x8D)  # Navy from template
        for j, col in enumerate(df_display.columns):
            cell = table.cell(0, j)
            cell.text = str(col)
            cell.fill.solid()
            cell.fill.fore_color.rgb = header_color
            p = cell.text_frame.paragraphs[0]
            p.font.bold = True
            p.font.size = Pt(9)
            p.font.color.rgb = RGBColor(255, 255, 255)
            p.font.name = 'Arial'
        
        # Data rows
        for i, row in df_display.iterrows():
            for j, val in enumerate(row):
                cell = table.cell(i + 1, j)
                cell.text = str(val) if pd.notna(val) else ""
                p = cell.text_frame.paragraphs[0]
                p.font.size = Pt(8)
                p.font.name = 'Arial'
        
        return slide
    
    def add_chart_slide(self, title: str, categories: List[str], series_data: Dict[str, List[float]], chart_type: str = "bar"):
        """Add a slide with a chart using template layout"""
        layout = self.layout_title_only or self.layout_blank or self.prs.slide_layouts[0]
        slide = self.prs.slides.add_slide(layout)
        
        # Set title via placeholder or fallback
        if self._using_template:
            self._set_placeholder_text(slide, 'title', title)
        else:
            title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12.333), Inches(0.6))
            title_frame = title_box.text_frame
            title_para = title_frame.paragraphs[0]
            title_para.text = title
            title_para.font.size = Pt(24)
            title_para.font.bold = True
        
        # Chart data
        chart_data = CategoryChartData()
        chart_data.categories = categories
        
        for series_name, values in series_data.items():
            chart_data.add_series(series_name, values)
        
        # Chart type
        if chart_type == "line":
            ct = XL_CHART_TYPE.LINE
        else:
            ct = XL_CHART_TYPE.COLUMN_CLUSTERED
        
        # Add chart with template-appropriate positioning
        chart_width = Inches(9.0) if self._using_template else Inches(12.333)
        x, y, cx, cy = Inches(0.4), Inches(1.2), chart_width, Inches(5.5)
        chart = slide.shapes.add_chart(ct, x, y, cx, cy, chart_data).chart
        
        return slide
    
    def add_kpi_slide(self, title: str, kpis: Dict[str, Any]):
        """Add a slide with KPI boxes using template layout"""
        layout = self.layout_title_only or self.layout_blank or self.prs.slide_layouts[0]
        slide = self.prs.slides.add_slide(layout)
        
        # Set title via placeholder or fallback
        if self._using_template:
            self._set_placeholder_text(slide, 'title', title)
        else:
            title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12.333), Inches(0.6))
        title_frame = title_box.text_frame
        title_para = title_frame.paragraphs[0]
        title_para.text = title
        title_para.font.size = Pt(24)
        title_para.font.bold = True
        
        # KPI boxes (up to 4)
        kpi_list = list(kpis.items())[:4]
        box_width = 2.8
        start_x = 0.7
        
        for i, (kpi_name, kpi_value) in enumerate(kpi_list):
            x = start_x + i * (box_width + 0.3)
            
            # Box
            shape = slide.shapes.add_shape(1, Inches(x), Inches(2), Inches(box_width), Inches(2))
            shape.fill.solid()
            shape.fill.fore_color.rgb = RGBColor(240, 240, 240)
            
            # KPI Value
            val_box = slide.shapes.add_textbox(Inches(x), Inches(2.3), Inches(box_width), Inches(1))
            val_frame = val_box.text_frame
            val_para = val_frame.paragraphs[0]
            val_para.text = str(kpi_value)
            val_para.font.size = Pt(32)
            val_para.font.bold = True
            val_para.alignment = 1  # Center
            
            # KPI Name
            name_box = slide.shapes.add_textbox(Inches(x), Inches(3.3), Inches(box_width), Inches(0.5))
            name_frame = name_box.text_frame
            name_para = name_frame.paragraphs[0]
            name_para.text = kpi_name
            name_para.font.size = Pt(14)
            name_para.alignment = 1  # Center
    
    def save(self):
        """Save the presentation"""
        self.prs.save(self.output_path)
        return self.output_path


class ExcelGenerator:
    """Generate Excel backup workbooks"""
    
    def __init__(self, output_path: Path):
        self.output_path = output_path
        self.wb = Workbook()
        self.wb.remove(self.wb.active)  # Remove default sheet
        
        # Styles
        self.header_font = Font(bold=True, color="FFFFFF")
        self.header_fill = PatternFill(start_color="2F5496", end_color="2F5496", fill_type="solid")
        self.border = Border(
            left=Side(style='thin'),
            right=Side(style='thin'),
            top=Side(style='thin'),
            bottom=Side(style='thin')
        )
    
    def add_dataframe_sheet(self, name: str, df: pd.DataFrame):
        """Add a sheet with a formatted dataframe"""
        ws = self.wb.create_sheet(title=name[:31])  # Excel 31 char limit
        
        # Headers
        for col_idx, column in enumerate(df.columns, 1):
            cell = ws.cell(row=1, column=col_idx, value=str(column))
            cell.font = self.header_font
            cell.fill = self.header_fill
            cell.border = self.border
            cell.alignment = Alignment(horizontal='center')
        
        # Data
        for row_idx, row in enumerate(df.itertuples(index=False), 2):
            for col_idx, value in enumerate(row, 1):
                cell = ws.cell(row=row_idx, column=col_idx, value=value)
                cell.border = self.border
        
        # Auto-width columns
        for col in ws.columns:
            max_length = max(len(str(cell.value or "")) for cell in col)
            ws.column_dimensions[col[0].column_letter].width = min(max_length + 2, 50)
    
    def add_summary_sheet(self, analyses: Dict[str, Any]):
        """Add a summary/overview sheet"""
        ws = self.wb.create_sheet(title="Summary", index=0)
        
        ws.cell(row=1, column=1, value="Data Pack Summary").font = Font(bold=True, size=16)
        ws.cell(row=2, column=1, value=f"Generated: {datetime.now().strftime('%Y-%m-%d %H:%M')}")
        
        row = 4
        for filename, sheets in analyses.items():
            ws.cell(row=row, column=1, value=filename).font = Font(bold=True)
            row += 1
            for sheet_name, data in sheets.items():
                ws.cell(row=row, column=2, value=sheet_name)
                row += 1
        
    def save(self):
        """Save the workbook"""
        self.wb.save(self.output_path)
        return self.output_path
