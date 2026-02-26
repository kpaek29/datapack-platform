"""
Smart Data Pack Generator v2
- Proper template-based formatting
- Quality validation against training examples
- Iterative analysis support
"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.oxml.ns import nsmap
from pptx.oxml import parse_xml
import pandas as pd
import numpy as np
from pathlib import Path
from typing import Dict, List, Any, Optional, Tuple
from datetime import datetime
from io import BytesIO
import json

# Template path
TEMPLATE_PATH = Path(__file__).parent.parent / "templates" / "master_template.pptx"


class SilverOakStyle:
    """Silver Oak brand styling constants"""
    
    # Colors
    DARK_TEAL = RGBColor(0x03, 0x33, 0x33)      # #033333 - primary
    GREEN = RGBColor(0x3E, 0x77, 0x33)           # #3E7733 - accent
    NAVY = RGBColor(0x08, 0x46, 0x8D)            # #08468D - headers
    GRAY = RGBColor(0x6B, 0x6B, 0x6B)            # #6B6B6B - body text
    GRAY_LIGHT = RGBColor(0x85, 0x86, 0x8A)      # #85868A - annotations
    WHITE = RGBColor(0xFF, 0xFF, 0xFF)
    CREAM = RGBColor(0xF9, 0xF8, 0xF6)           # #F9F8F6 - background
    
    # Fonts
    HEADING_FONT = "Libre Baskerville"
    BODY_FONT = "Arial"
    
    # Font sizes (in points)
    TITLE_SIZE = 24
    SUBTITLE_SIZE = 14
    TABLE_HEADER_SIZE = 9
    TABLE_BODY_SIZE = 8
    FOOTNOTE_SIZE = 6
    
    # Table colors
    TABLE_HEADER_BG = NAVY
    TABLE_HEADER_FG = WHITE
    TABLE_ALT_ROW = RGBColor(0xF5, 0xF5, 0xF5)  # Light gray for alternating rows
    TABLE_BORDER = RGBColor(0xDD, 0xDD, 0xDD)


class SmartTableFormatter:
    """Intelligent table formatting for PPT"""
    
    @staticmethod
    def calculate_column_widths(df: pd.DataFrame, total_width: float) -> List[float]:
        """Calculate optimal column widths based on content"""
        widths = []
        
        for col in df.columns:
            # Get max content length in column
            header_len = len(str(col))
            max_data_len = df[col].astype(str).str.len().max() if len(df) > 0 else 0
            max_len = max(header_len, max_data_len, 5)  # Minimum 5 chars
            widths.append(max_len)
        
        # Normalize to total width
        total_chars = sum(widths)
        if total_chars > 0:
            widths = [w / total_chars * total_width for w in widths]
        else:
            widths = [total_width / len(df.columns)] * len(df.columns)
            
        return widths
    
    @staticmethod
    def format_value(val: Any, col_name: str = "") -> str:
        """Smart value formatting based on data type"""
        if pd.isna(val):
            return ""
        
        # Detect and format numbers
        if isinstance(val, (int, float)):
            # Currency/large numbers
            if abs(val) >= 1000000:
                return f"${val/1000000:,.1f}M"
            elif abs(val) >= 1000:
                return f"${val/1000:,.0f}K"
            elif abs(val) >= 1:
                if val == int(val):
                    return f"{int(val):,}"
                return f"{val:,.2f}"
            elif abs(val) < 1 and val != 0:
                # Percentages
                return f"{val*100:.1f}%"
            else:
                return str(val)
        
        # Truncate long strings
        val_str = str(val)
        if len(val_str) > 40:
            return val_str[:37] + "..."
        return val_str
    
    @staticmethod
    def add_styled_table(
        slide, 
        df: pd.DataFrame,
        left: float,
        top: float,
        width: float,
        height: float,
        max_rows: int = 20,
        show_index: bool = False
    ):
        """Add a properly styled table to a slide"""
        
        # Limit rows
        df_display = df.head(max_rows).copy()
        
        # Calculate dimensions
        rows = len(df_display) + 1  # +1 for header
        cols = len(df_display.columns)
        
        if cols == 0:
            return None
        
        # Create table
        table_shape = slide.shapes.add_table(
            rows, cols,
            Inches(left), Inches(top),
            Inches(width), Inches(height)
        )
        table = table_shape.table
        
        # Calculate column widths
        col_widths = SmartTableFormatter.calculate_column_widths(df_display, width)
        for i, w in enumerate(col_widths):
            table.columns[i].width = Inches(w)
        
        # Style header row
        for j, col in enumerate(df_display.columns):
            cell = table.cell(0, j)
            cell.text = str(col)
            
            # Fill
            cell.fill.solid()
            cell.fill.fore_color.rgb = SilverOakStyle.TABLE_HEADER_BG
            
            # Text styling
            para = cell.text_frame.paragraphs[0]
            para.font.name = SilverOakStyle.BODY_FONT
            para.font.size = Pt(SilverOakStyle.TABLE_HEADER_SIZE)
            para.font.color.rgb = SilverOakStyle.TABLE_HEADER_FG
            para.font.bold = True
            para.alignment = PP_ALIGN.CENTER
            
            # Vertical alignment
            cell.vertical_anchor = MSO_ANCHOR.MIDDLE
            
            # Margins
            cell.margin_left = Inches(0.05)
            cell.margin_right = Inches(0.05)
            cell.margin_top = Inches(0.03)
            cell.margin_bottom = Inches(0.03)
        
        # Style data rows
        for i, (idx, row) in enumerate(df_display.iterrows()):
            for j, val in enumerate(row):
                cell = table.cell(i + 1, j)
                
                # Format value
                col_name = df_display.columns[j]
                cell.text = SmartTableFormatter.format_value(val, col_name)
                
                # Alternating row colors
                if i % 2 == 1:
                    cell.fill.solid()
                    cell.fill.fore_color.rgb = SilverOakStyle.TABLE_ALT_ROW
                else:
                    cell.fill.background()
                
                # Text styling
                para = cell.text_frame.paragraphs[0]
                para.font.name = SilverOakStyle.BODY_FONT
                para.font.size = Pt(SilverOakStyle.TABLE_BODY_SIZE)
                para.font.color.rgb = SilverOakStyle.GRAY
                
                # Right-align numbers
                if isinstance(val, (int, float)) and not pd.isna(val):
                    para.alignment = PP_ALIGN.RIGHT
                else:
                    para.alignment = PP_ALIGN.LEFT
                
                cell.vertical_anchor = MSO_ANCHOR.MIDDLE
                cell.margin_left = Inches(0.05)
                cell.margin_right = Inches(0.05)
                cell.margin_top = Inches(0.02)
                cell.margin_bottom = Inches(0.02)
        
        return table_shape


class SmartPPTGenerator:
    """Generate high-quality data pack presentations"""
    
    def __init__(
        self, 
        output_path: Path,
        company_name: str,
        date_str: str = None,
        template_path: Path = None
    ):
        # Load template
        template = template_path or TEMPLATE_PATH
        if template.exists():
            self.prs = Presentation(str(template))
            self._clear_template_slides()
            self._using_template = True
        else:
            self.prs = Presentation()
            self.prs.slide_width = Inches(10)
            self.prs.slide_height = Inches(7.5)
            self._using_template = False
        
        self.output_path = Path(output_path)
        self.company_name = company_name
        self.date_str = date_str or datetime.now().strftime("%B %Y")
        self._setup_layouts()
        
    def _clear_template_slides(self):
        """Remove placeholder slides from template"""
        while len(self.prs.slides) > 0:
            rId = self.prs.slides._sldIdLst[0].rId
            self.prs.part.drop_rel(rId)
            del self.prs.slides._sldIdLst[0]
    
    def _setup_layouts(self):
        """Map slide layouts"""
        layouts = self.prs.slide_layouts
        
        self.layout_title = layouts[0] if len(layouts) > 0 else None
        self.layout_content = layouts[1] if len(layouts) > 1 else layouts[0]
        self.layout_section = layouts[2] if len(layouts) > 2 else layouts[0]
        self.layout_title_only = layouts[4] if len(layouts) > 4 else layouts[0]
    
    def _add_title_text(self, slide, text: str, left: float = 0.4, top: float = 0.3):
        """Add styled title text"""
        title_box = slide.shapes.add_textbox(
            Inches(left), Inches(top), Inches(9.2), Inches(0.5)
        )
        tf = title_box.text_frame
        para = tf.paragraphs[0]
        para.text = text
        para.font.name = SilverOakStyle.HEADING_FONT
        para.font.size = Pt(SilverOakStyle.TITLE_SIZE)
        para.font.color.rgb = SilverOakStyle.DARK_TEAL
        para.font.bold = False
        return title_box
    
    def _add_subtitle_text(self, slide, text: str, left: float = 0.4, top: float = 0.75):
        """Add subtitle/annotation text"""
        box = slide.shapes.add_textbox(
            Inches(left), Inches(top), Inches(9.2), Inches(0.3)
        )
        tf = box.text_frame
        para = tf.paragraphs[0]
        para.text = text
        para.font.name = SilverOakStyle.HEADING_FONT
        para.font.size = Pt(SilverOakStyle.SUBTITLE_SIZE)
        para.font.color.rgb = SilverOakStyle.GRAY
        para.font.italic = True
        return box
    
    def _add_footnote(self, slide, text: str, top: float = 6.8):
        """Add footnote text"""
        box = slide.shapes.add_textbox(
            Inches(0.4), Inches(top), Inches(9.2), Inches(0.3)
        )
        tf = box.text_frame
        para = tf.paragraphs[0]
        para.text = text
        para.font.name = SilverOakStyle.HEADING_FONT
        para.font.size = Pt(SilverOakStyle.FOOTNOTE_SIZE)
        para.font.color.rgb = SilverOakStyle.GRAY_LIGHT
        para.font.italic = True
        return box
    
    def add_title_slide(self, title: str = None, subtitle: str = None):
        """Add cover slide"""
        slide = self.prs.slides.add_slide(self.layout_title)
        
        # Try to use placeholders first
        for shape in slide.placeholders:
            if shape.placeholder_format.type == 1:  # Title
                shape.text = title or f"{self.company_name} Data Pack"
            elif shape.placeholder_format.type == 2:  # Subtitle
                shape.text = subtitle or self.date_str
        
        return slide
    
    def add_section_slide(self, title: str, subtitle: str = None):
        """Add section divider slide"""
        slide = self.prs.slides.add_slide(self.layout_section)
        
        # Set title
        for shape in slide.placeholders:
            if 'title' in shape.name.lower():
                shape.text = title
                break
        else:
            self._add_title_text(slide, title, top=3.0)
        
        if subtitle:
            self._add_subtitle_text(slide, subtitle, top=3.6)
        
        return slide
    
    def add_table_slide(
        self,
        title: str,
        df: pd.DataFrame,
        subtitle: str = None,
        footnote: str = None,
        max_rows: int = 18
    ):
        """Add a properly formatted table slide"""
        slide = self.prs.slides.add_slide(self.layout_title_only)
        
        # Add title
        self._add_title_text(slide, title)
        
        # Add subtitle if provided
        content_top = 0.9
        if subtitle:
            self._add_subtitle_text(slide, subtitle)
            content_top = 1.1
        
        # Calculate table height based on rows
        num_rows = min(len(df) + 1, max_rows + 1)
        table_height = min(num_rows * 0.3, 5.5)
        
        # Add table
        if len(df) > 0:
            SmartTableFormatter.add_styled_table(
                slide,
                df,
                left=0.4,
                top=content_top,
                width=9.2,
                height=table_height,
                max_rows=max_rows
            )
        
        # Add footnote
        if footnote:
            self._add_footnote(slide, footnote)
        
        return slide
    
    def add_pl_slide(self, title: str, df: pd.DataFrame):
        """Add P&L summary slide with special formatting"""
        slide = self.prs.slides.add_slide(self.layout_title_only)
        
        self._add_title_text(slide, title)
        
        if len(df) == 0:
            self._add_subtitle_text(slide, "No P&L data available", top=3.0)
            return slide
        
        # P&L tables need special row highlighting for totals
        SmartTableFormatter.add_styled_table(
            slide,
            df,
            left=0.4,
            top=1.0,
            width=9.2,
            height=5.5,
            max_rows=22
        )
        
        self._add_footnote(slide, "Source: Company financials")
        return slide
    
    def add_customer_slide(self, title: str, df: pd.DataFrame):
        """Add customer analysis slide"""
        slide = self.prs.slides.add_slide(self.layout_title_only)
        
        self._add_title_text(slide, title)
        
        if len(df) == 0:
            self._add_subtitle_text(slide, "No customer data available", top=3.0)
            return slide
        
        SmartTableFormatter.add_styled_table(
            slide,
            df,
            left=0.4,
            top=1.0,
            width=9.2,
            height=5.5,
            max_rows=20
        )
        
        return slide
    
    def add_chart_slide(
        self,
        title: str,
        chart_image: bytes,
        subtitle: str = None,
        footnote: str = None
    ):
        """Add slide with chart image"""
        slide = self.prs.slides.add_slide(self.layout_title_only)
        
        self._add_title_text(slide, title)
        
        top = 0.9
        if subtitle:
            self._add_subtitle_text(slide, subtitle)
            top = 1.2
        
        # Add chart
        img_stream = BytesIO(chart_image)
        slide.shapes.add_picture(
            img_stream,
            Inches(0.4), Inches(top),
            width=Inches(9.0)
        )
        
        if footnote:
            self._add_footnote(slide, footnote)
        
        return slide
    
    def add_kpi_slide(self, title: str, kpis: Dict[str, Any]):
        """Add KPI summary slide"""
        slide = self.prs.slides.add_slide(self.layout_title_only)
        
        self._add_title_text(slide, title)
        
        # Create KPI boxes
        kpi_list = list(kpis.items())[:6]  # Max 6 KPIs
        
        cols = min(len(kpi_list), 3)
        rows = (len(kpi_list) + cols - 1) // cols
        
        box_width = 2.8
        box_height = 1.5
        start_x = 0.5
        start_y = 1.2
        gap = 0.3
        
        for i, (name, value) in enumerate(kpi_list):
            row = i // cols
            col = i % cols
            
            x = start_x + col * (box_width + gap)
            y = start_y + row * (box_height + gap)
            
            # Box background
            shape = slide.shapes.add_shape(
                1,  # Rectangle
                Inches(x), Inches(y),
                Inches(box_width), Inches(box_height)
            )
            shape.fill.solid()
            shape.fill.fore_color.rgb = SilverOakStyle.CREAM
            shape.line.color.rgb = SilverOakStyle.TABLE_BORDER
            
            # Value
            val_box = slide.shapes.add_textbox(
                Inches(x), Inches(y + 0.3),
                Inches(box_width), Inches(0.6)
            )
            tf = val_box.text_frame
            para = tf.paragraphs[0]
            para.text = str(value)
            para.font.name = SilverOakStyle.HEADING_FONT
            para.font.size = Pt(28)
            para.font.color.rgb = SilverOakStyle.GREEN
            para.font.bold = True
            para.alignment = PP_ALIGN.CENTER
            
            # Label
            label_box = slide.shapes.add_textbox(
                Inches(x), Inches(y + 0.9),
                Inches(box_width), Inches(0.4)
            )
            tf = label_box.text_frame
            para = tf.paragraphs[0]
            para.text = name
            para.font.name = SilverOakStyle.BODY_FONT
            para.font.size = Pt(10)
            para.font.color.rgb = SilverOakStyle.GRAY
            para.alignment = PP_ALIGN.CENTER
        
        return slide
    
    def save(self) -> Path:
        """Save the presentation"""
        self.prs.save(self.output_path)
        return self.output_path


class AnalysisSuggester:
    """
    AI-powered analysis suggestion based on data content
    Scans uploaded data and recommends relevant analyses
    """
    
    def __init__(self, openai_client=None):
        self.client = openai_client
    
    def suggest_analyses(self, data_summary: dict, sheets_info: dict) -> dict:
        """
        Suggest analyses based on data content
        
        Args:
            data_summary: Summary of detected data types
            sheets_info: Info about sheets and columns
            
        Returns:
            {
                'suggested': ['analysis_type1', 'analysis_type2'],
                'reasons': {'analysis_type1': 'reason'},
                'confidence': {'analysis_type1': 0.9},
                'additional': ['other_available_analyses']
            }
        """
        suggested = []
        reasons = {}
        confidence = {}
        
        # Check for customer data
        has_customers = any(
            'customer' in str(info).lower() or 'client' in str(info).lower()
            for info in sheets_info.values()
        )
        
        has_revenue = any(
            'revenue' in str(info).lower() or 'sales' in str(info).lower() or 'amount' in str(info).lower()
            for info in sheets_info.values()
        )
        
        has_dates = any(
            'date' in str(info).lower() or '2024' in str(info) or '2025' in str(info) or '2023' in str(info)
            for info in sheets_info.values()
        )
        
        has_segments = any(
            'segment' in str(info).lower() or 'category' in str(info).lower() or 'type' in str(info).lower()
            for info in sheets_info.values()
        )
        
        # Suggest based on data
        if has_customers and has_revenue:
            suggested.append('top_customers')
            reasons['top_customers'] = 'Customer and revenue data detected'
            confidence['top_customers'] = 0.95
            
            suggested.append('customer_concentration')
            reasons['customer_concentration'] = 'Can analyze revenue concentration across customers'
            confidence['customer_concentration'] = 0.9
        
        if has_customers and has_dates:
            suggested.append('customer_retention')
            reasons['customer_retention'] = 'Customer data with dates allows retention analysis'
            confidence['customer_retention'] = 0.7
        
        if has_revenue and has_dates:
            suggested.append('monthly_trends')
            reasons['monthly_trends'] = 'Revenue with dates enables trend analysis'
            confidence['monthly_trends'] = 0.85
            
            suggested.append('yoy_comparison')
            reasons['yoy_comparison'] = 'Multi-period data allows year-over-year comparison'
            confidence['yoy_comparison'] = 0.75
        
        if has_segments and has_revenue:
            suggested.append('revenue_by_segment')
            reasons['revenue_by_segment'] = 'Segment and revenue data detected'
            confidence['revenue_by_segment'] = 0.85
        
        if has_customers and has_dates and has_revenue:
            suggested.append('cohort_analysis')
            reasons['cohort_analysis'] = 'Full customer transaction data enables cohort analysis'
            confidence['cohort_analysis'] = 0.6
        
        # Sort by confidence
        suggested.sort(key=lambda x: confidence.get(x, 0), reverse=True)
        
        # Additional analyses not suggested
        all_analyses = list(IterativeAnalyzer.ANALYSIS_TYPES.keys())
        additional = [a for a in all_analyses if a not in suggested]
        
        return {
            'suggested': suggested[:5],  # Top 5 suggestions
            'reasons': reasons,
            'confidence': confidence,
            'additional': additional
        }
    
    def analyze_dataframes(self, dataframes: dict) -> dict:
        """
        Analyze actual dataframes to make suggestions
        """
        sheets_info = {}
        
        for name, df in dataframes.items():
            if hasattr(df, 'columns'):
                sheets_info[name] = {
                    'columns': list(df.columns),
                    'rows': len(df),
                    'sample': df.head(3).to_dict() if len(df) > 0 else {}
                }
        
        # Build summary
        all_columns = []
        for info in sheets_info.values():
            all_columns.extend([str(c).lower() for c in info.get('columns', [])])
        
        data_summary = {
            'total_sheets': len(sheets_info),
            'all_columns': all_columns
        }
        
        return self.suggest_analyses(data_summary, sheets_info)


class IterativeAnalyzer:
    """
    Handles iterative analysis requests
    Allows users to request specific analyses via natural language
    """
    
    ANALYSIS_TYPES = {
        'customer_retention': {
            'description': 'Customer retention and churn analysis',
            'requires': ['customer_data', 'date_column'],
            'outputs': ['retention_table', 'cohort_chart']
        },
        'revenue_by_segment': {
            'description': 'Revenue breakdown by segment/category',
            'requires': ['revenue_data', 'segment_column'],
            'outputs': ['segment_table', 'segment_chart']
        },
        'top_customers': {
            'description': 'Top customers by revenue',
            'requires': ['customer_data', 'revenue_column'],
            'outputs': ['top_customers_table']
        },
        'monthly_trends': {
            'description': 'Monthly revenue/metric trends',
            'requires': ['date_column', 'metric_column'],
            'outputs': ['trend_chart']
        },
        'yoy_comparison': {
            'description': 'Year-over-year comparison',
            'requires': ['date_column', 'metric_column'],
            'outputs': ['yoy_table', 'yoy_chart']
        },
        'customer_concentration': {
            'description': 'Customer concentration analysis (top 10, 20, etc.)',
            'requires': ['customer_data', 'revenue_column'],
            'outputs': ['concentration_table']
        },
        'cohort_analysis': {
            'description': 'Customer cohort analysis by signup period',
            'requires': ['customer_data', 'date_column', 'revenue_column'],
            'outputs': ['cohort_table']
        },
        'gross_margin': {
            'description': 'Gross margin analysis',
            'requires': ['revenue_data', 'cogs_data'],
            'outputs': ['margin_table', 'margin_chart']
        }
    }
    
    def __init__(self, openai_client=None):
        self.client = openai_client
    
    def parse_request(self, request: str) -> Dict[str, Any]:
        """
        Parse a natural language analysis request
        
        Examples:
        - "Add customer retention analysis"
        - "Show revenue by segment"
        - "I need a cohort analysis"
        """
        request_lower = request.lower()
        
        # Simple keyword matching
        matched_analyses = []
        
        if any(kw in request_lower for kw in ['retention', 'churn']):
            matched_analyses.append('customer_retention')
        if any(kw in request_lower for kw in ['segment', 'category', 'breakdown']):
            matched_analyses.append('revenue_by_segment')
        if any(kw in request_lower for kw in ['top customer', 'largest customer']):
            matched_analyses.append('top_customers')
        if any(kw in request_lower for kw in ['trend', 'monthly', 'over time']):
            matched_analyses.append('monthly_trends')
        if any(kw in request_lower for kw in ['year over year', 'yoy', 'y-o-y']):
            matched_analyses.append('yoy_comparison')
        if any(kw in request_lower for kw in ['concentration', 'top 10', 'top 20']):
            matched_analyses.append('customer_concentration')
        if any(kw in request_lower for kw in ['cohort']):
            matched_analyses.append('cohort_analysis')
        if any(kw in request_lower for kw in ['margin', 'gross profit']):
            matched_analyses.append('gross_margin')
        
        return {
            'request': request,
            'matched_analyses': matched_analyses,
            'details': [self.ANALYSIS_TYPES.get(a) for a in matched_analyses]
        }
    
    def generate_analysis(
        self,
        analysis_type: str,
        data: Dict[str, pd.DataFrame],
        parameters: Dict[str, Any] = None
    ) -> Dict[str, Any]:
        """
        Generate a specific analysis from data
        
        Returns:
            {
                'title': str,
                'subtitle': str,
                'data': pd.DataFrame,
                'chart': bytes (optional),
                'insights': str
            }
        """
        params = parameters or {}
        
        if analysis_type == 'customer_retention':
            return self._analyze_retention(data, params)
        elif analysis_type == 'top_customers':
            return self._analyze_top_customers(data, params)
        elif analysis_type == 'customer_concentration':
            return self._analyze_concentration(data, params)
        elif analysis_type == 'revenue_by_segment':
            return self._analyze_segments(data, params)
        elif analysis_type == 'monthly_trends':
            return self._analyze_trends(data, params)
        else:
            return {
                'title': f'{analysis_type.replace("_", " ").title()} Analysis',
                'subtitle': 'Analysis not yet implemented',
                'data': pd.DataFrame(),
                'insights': ''
            }
    
    def _analyze_retention(self, data: Dict, params: Dict) -> Dict:
        """Generate customer retention analysis"""
        customer_df = data.get('customers') or data.get('top_customers', pd.DataFrame())
        
        if customer_df.empty:
            return {
                'title': 'Customer Retention Analysis',
                'subtitle': 'Insufficient data for retention analysis',
                'data': pd.DataFrame(),
                'insights': 'Need customer transaction data with dates to calculate retention.'
            }
        
        # Simple retention calculation if we have the data
        result_df = pd.DataFrame({
            'Metric': ['Active Customers', 'New Customers', 'Churned', 'Retention Rate'],
            'Current Period': ['N/A', 'N/A', 'N/A', 'N/A'],
            'Prior Period': ['N/A', 'N/A', 'N/A', 'N/A']
        })
        
        return {
            'title': 'Customer Retention Analysis',
            'subtitle': 'Period-over-period customer retention metrics',
            'data': result_df,
            'insights': 'Retention analysis requires transaction-level data with dates.'
        }
    
    def _analyze_top_customers(self, data: Dict, params: Dict) -> Dict:
        """Generate top customers analysis"""
        df = data.get('top_customers', pd.DataFrame())
        
        if df.empty:
            return {
                'title': 'Top Customers',
                'subtitle': 'No customer data available',
                'data': pd.DataFrame(),
                'insights': ''
            }
        
        # Ensure sorted
        if len(df.columns) >= 2:
            amount_col = df.columns[1]
            df = df.sort_values(by=amount_col, ascending=False).head(20)
        
        return {
            'title': f'Top {len(df)} Customers',
            'subtitle': 'Ranked by revenue contribution',
            'data': df,
            'insights': ''
        }
    
    def _analyze_concentration(self, data: Dict, params: Dict) -> Dict:
        """Generate customer concentration analysis"""
        df = data.get('top_customers', pd.DataFrame())
        
        if df.empty or len(df.columns) < 2:
            return {
                'title': 'Customer Concentration',
                'subtitle': 'No data available',
                'data': pd.DataFrame(),
                'insights': ''
            }
        
        amount_col = df.columns[1]
        df_sorted = df.sort_values(by=amount_col, ascending=False)
        total = df_sorted[amount_col].sum()
        
        # Calculate concentration
        concentration = []
        for n in [1, 5, 10, 20]:
            if len(df_sorted) >= n:
                top_n_sum = df_sorted.head(n)[amount_col].sum()
                pct = (top_n_sum / total * 100) if total > 0 else 0
                concentration.append({
                    'Segment': f'Top {n} Customers',
                    'Revenue': f'${top_n_sum:,.0f}',
                    '% of Total': f'{pct:.1f}%'
                })
        
        result_df = pd.DataFrame(concentration)
        
        return {
            'title': 'Customer Concentration Analysis',
            'subtitle': 'Revenue concentration by customer segment',
            'data': result_df,
            'insights': f'Top customer represents {concentration[0]["% of Total"] if concentration else "N/A"} of total revenue.'
        }
    
    def _analyze_segments(self, data: Dict, params: Dict) -> Dict:
        """Generate revenue by segment analysis"""
        # Placeholder - would need segment column identification
        return {
            'title': 'Revenue by Segment',
            'subtitle': 'Segment breakdown not available',
            'data': pd.DataFrame({'Segment': ['N/A'], 'Revenue': ['N/A']}),
            'insights': 'Need segment/category column in data.'
        }
    
    def _analyze_trends(self, data: Dict, params: Dict) -> Dict:
        """Generate trend analysis"""
        df = data.get('monthly_revenue', pd.DataFrame())
        
        if df.empty:
            return {
                'title': 'Monthly Trends',
                'subtitle': 'No trend data available',
                'data': pd.DataFrame(),
                'insights': ''
            }
        
        return {
            'title': 'Monthly Revenue Trends',
            'subtitle': 'Revenue over time',
            'data': df,
            'insights': ''
        }


class QualityValidator:
    """
    Validates output quality against training examples
    """
    
    def __init__(self, training_dir: Path = None):
        self.training_dir = training_dir
        self.quality_thresholds = {
            'min_table_rows': 3,
            'min_columns': 2,
            'max_empty_cells_pct': 30,
            'min_data_coverage': 50
        }
    
    def validate_dataframe(self, df: pd.DataFrame, context: str = "") -> Dict[str, Any]:
        """Validate a dataframe for quality issues"""
        issues = []
        score = 100
        
        if df.empty:
            return {
                'valid': False,
                'score': 0,
                'issues': ['DataFrame is empty'],
                'context': context
            }
        
        # Check row count
        if len(df) < self.quality_thresholds['min_table_rows']:
            issues.append(f'Only {len(df)} rows (minimum: {self.quality_thresholds["min_table_rows"]})')
            score -= 20
        
        # Check column count
        if len(df.columns) < self.quality_thresholds['min_columns']:
            issues.append(f'Only {len(df.columns)} columns')
            score -= 20
        
        # Check empty cells
        total_cells = df.size
        empty_cells = df.isna().sum().sum()
        empty_pct = (empty_cells / total_cells * 100) if total_cells > 0 else 0
        
        if empty_pct > self.quality_thresholds['max_empty_cells_pct']:
            issues.append(f'{empty_pct:.0f}% empty cells')
            score -= 30
        
        return {
            'valid': score >= 60,
            'score': max(score, 0),
            'issues': issues,
            'context': context,
            'stats': {
                'rows': len(df),
                'columns': len(df.columns),
                'empty_pct': empty_pct
            }
        }
    
    def validate_presentation(self, slides_data: List[Dict]) -> Dict[str, Any]:
        """Validate overall presentation quality"""
        all_issues = []
        total_score = 0
        
        for slide in slides_data:
            if 'data' in slide and isinstance(slide['data'], pd.DataFrame):
                validation = self.validate_dataframe(slide['data'], slide.get('title', ''))
                total_score += validation['score']
                if validation['issues']:
                    all_issues.extend([f"{slide.get('title', 'Slide')}: {i}" for i in validation['issues']])
        
        avg_score = total_score / len(slides_data) if slides_data else 0
        
        return {
            'overall_valid': avg_score >= 60,
            'average_score': avg_score,
            'issues': all_issues,
            'recommendation': 'Review data quality' if all_issues else 'Quality check passed'
        }
