"""
DataPack Generator - PE Data Pack Creation Engine
Generates PPT and Excel outputs using master slide template
"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
import pandas as pd
import matplotlib.pyplot as plt
import matplotlib
matplotlib.use('Agg')
from io import BytesIO
from pathlib import Path
from datetime import datetime
from typing import Dict, List, Optional
import tempfile

# ============ TEMPLATE & STYLING CONSTANTS ============

# Path to master slide template
TEMPLATE_PATH = Path(__file__).parent.parent / "templates" / "master_template.pptx"

class DataPackStyle:
    """Silver Oak PE Data Pack styling - derived from master template"""
    
    # Colors from template
    DARK = RGBColor(0x51, 0x51, 0x51)      # #515151 - charcoal
    LIGHT = RGBColor(0xE5, 0xE5, 0xE5)     # #E5E5E5 - light gray
    GREEN = RGBColor(0x3E, 0x77, 0x33)      # #3E7733 - accent green
    NAVY = RGBColor(0x08, 0x46, 0x8D)       # #08468D - accent blue
    WHITE = RGBColor(0xFF, 0xFF, 0xFF)
    GRAY_TEXT = RGBColor(0x6B, 0x6B, 0x6B)  # #6B6B6B - gray text
    ANNOTATION_GRAY = RGBColor(0x85, 0x86, 0x8A)  # #85868A - annotation gray
    
    # Matplotlib colors
    MPL_DARK = '#515151'
    MPL_GREEN = '#3E7733'
    MPL_NAVY = '#08468D'
    MPL_LIGHT = '#E5E5E5'
    
    # Fonts from template
    HEADING_FONT = 'Libre Baskerville'  # Template uses Libre Baskerville
    BODY_FONT = 'Arial'
    ANNOTATION_FONT = 'Libre Baskerville Italic'
    
    # Template slide layout indices (from master_template.pptx)
    LAYOUT_TITLE = 0        # Title Slide
    LAYOUT_CONTENT = 1      # Title and Content
    LAYOUT_SECTION = 2      # Section Header
    LAYOUT_TWO_CONTENT = 3  # Two Content
    LAYOUT_TITLE_ONLY = 4   # Title Only


class DataPackPPTGenerator:
    """Generate PE Data Pack PowerPoint presentations using master template"""
    
    def __init__(self, output_path: str, company_name: str, date_str: str = None, 
                 template_path: Path = None):
        # Load from master template
        template = template_path or TEMPLATE_PATH
        if template.exists():
            self.prs = Presentation(str(template))
            # Remove placeholder slides from template
            while len(self.prs.slides) > 0:
                rId = self.prs.slides._sldIdLst[0].rId
                self.prs.part.drop_rel(rId)
                del self.prs.slides._sldIdLst[0]
        else:
            # Fallback to blank if template missing
            self.prs = Presentation()
            self.prs.slide_width = Inches(10)
            self.prs.slide_height = Inches(7.5)
            
        self.output_path = Path(output_path)
        self.company_name = company_name
        self.date_str = date_str or datetime.now().strftime("%B %Y")
        self.page_num = 0
        
        # Use layouts from template
        self._setup_layouts()
    
    def _setup_layouts(self):
        """Setup slide layout references from template"""
        layouts = self.prs.slide_layouts
        
        # Map layouts by name or index
        self.layout_title = None
        self.layout_content = None
        self.layout_section = None
        self.layout_two_content = None
        self.layout_title_only = None
        
        for i, layout in enumerate(layouts):
            name = layout.name.lower()
            if 'title slide' in name or (i == 0 and 'title' in name):
                self.layout_title = layout
            elif 'title and content' in name or 'content' in name.replace('two', ''):
                if self.layout_content is None:
                    self.layout_content = layout
            elif 'section' in name:
                self.layout_section = layout
            elif 'two content' in name:
                self.layout_two_content = layout
            elif 'title only' in name:
                self.layout_title_only = layout
        
        # Fallback to index-based if names don't match
        if len(layouts) >= 5:
            self.layout_title = self.layout_title or layouts[0]
            self.layout_content = self.layout_content or layouts[1]
            self.layout_section = self.layout_section or layouts[2]
            self.layout_two_content = self.layout_two_content or layouts[3]
            self.layout_title_only = self.layout_title_only or layouts[4]
        elif len(layouts) > 0:
            # Use first layout as fallback
            fallback = layouts[0]
            self.layout_title = self.layout_title or fallback
            self.layout_content = self.layout_content or fallback
            self.layout_section = self.layout_section or fallback
            self.layout_two_content = self.layout_two_content or fallback
            self.layout_title_only = self.layout_title_only or fallback
    
    def _set_placeholder_text(self, slide, ph_type: str, text: str):
        """Set text in a placeholder by type (title, body, subtitle)"""
        from pptx.enum.shapes import PP_PLACEHOLDER
        
        ph_map = {
            'title': PP_PLACEHOLDER.TITLE,
            'body': PP_PLACEHOLDER.BODY,
            'subtitle': PP_PLACEHOLDER.SUBTITLE,
            'center_title': PP_PLACEHOLDER.CENTER_TITLE,
        }
        
        target_type = ph_map.get(ph_type)
        
        for shape in slide.placeholders:
            if target_type and shape.placeholder_format.type == target_type:
                shape.text = text
                return shape
            elif ph_type == 'title' and 'title' in shape.name.lower():
                shape.text = text
                return shape
            elif ph_type == 'body' and ('content' in shape.name.lower() or 'body' in shape.name.lower()):
                shape.text = text
                return shape
        return None
    
    def _add_text_box(self, slide, text: str, left: float, top: float, 
                      width: float = 9.0, height: float = 0.4,
                      font_size: int = 12, bold: bool = False, 
                      font_name: str = None, color: RGBColor = None):
        """Add a text box with styling"""
        box = slide.shapes.add_textbox(
            Inches(left), Inches(top), Inches(width), Inches(height)
        )
        tf = box.text_frame
        p = tf.paragraphs[0]
        p.text = text
        p.font.name = font_name or DataPackStyle.BODY_FONT
        p.font.size = Pt(font_size)
        p.font.color.rgb = color or DataPackStyle.DARK
        p.font.bold = bold
        return box
        
    def add_title_slide(self):
        """Add cover/title slide using template layout"""
        slide = self.prs.slides.add_slide(self.layout_title)
        
        # Set title and subtitle placeholders
        self._set_placeholder_text(slide, 'title', f"{self.company_name} Data Pack")
        self._set_placeholder_text(slide, 'subtitle', self.date_str)
        
        return slide
        
    def add_agenda_slide(self, sections: List[str], current_section: str = None):
        """Add agenda slide using content layout"""
        slide = self.prs.slides.add_slide(self.layout_content)
        
        # Set title
        self._set_placeholder_text(slide, 'title', "Agenda")
        
        # Build agenda text with bullet points
        agenda_text = "\n".join(sections)
        body_shape = self._set_placeholder_text(slide, 'body', "")
        
        if body_shape:
            tf = body_shape.text_frame
            tf.clear()
            for i, section in enumerate(sections):
                p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
                p.text = section
                p.level = 0
                if section == current_section:
                    p.font.bold = True
        else:
            # Fallback: add text box
            y_pos = 1.2
            for section in sections:
                is_current = section == current_section
                self._add_text_box(
                    slide, section, 0.6, y_pos,
                    font_size=14,
                    bold=is_current
                )
                y_pos += 0.4
        
        return slide
            
    def add_section_slide(self, title: str, subtitle: str = None):
        """Add section header slide using template layout"""
        slide = self.prs.slides.add_slide(self.layout_section)
        
        self._set_placeholder_text(slide, 'title', title)
        if subtitle:
            self._set_placeholder_text(slide, 'body', subtitle)
        
        return slide
        
    def add_pl_summary_slide(self, title: str, data: pd.DataFrame):
        """Add P&L summary slide with table using content layout"""
        slide = self.prs.slides.add_slide(self.layout_title_only)
        
        # Set title using placeholder
        self._set_placeholder_text(slide, 'title', title)
        
        # Create table
        rows = min(len(data) + 1, 20)  # Header + data, max 20
        cols = min(len(data.columns), 10)
        
        # Position table below title (respecting template spacing)
        table = slide.shapes.add_table(
            rows, cols,
            Inches(0.4), Inches(1.2),
            Inches(9.2), Inches(5.5)
        ).table
        
        # Style header row
        for j, col in enumerate(data.columns[:cols]):
            cell = table.cell(0, j)
            cell.text = str(col)
            cell.fill.solid()
            cell.fill.fore_color.rgb = DataPackStyle.NAVY
            p = cell.text_frame.paragraphs[0]
            p.font.color.rgb = DataPackStyle.WHITE
            p.font.size = Pt(9)
            p.font.bold = True
            p.font.name = DataPackStyle.BODY_FONT
            
        # Fill data
        for i, (idx, row) in enumerate(data.iterrows()):
            if i >= rows - 1:
                break
            for j, val in enumerate(row[:cols]):
                cell = table.cell(i + 1, j)
                cell.text = str(val) if pd.notna(val) else ""
                p = cell.text_frame.paragraphs[0]
                p.font.size = Pt(8)
                p.font.name = DataPackStyle.BODY_FONT
        
        return slide
                
    def add_chart_slide(self, title: str, chart_image: bytes, 
                        subtitle: str = None, footnote: str = None):
        """Add slide with chart image using title only layout"""
        slide = self.prs.slides.add_slide(self.layout_title_only)
        
        # Set title using placeholder
        self._set_placeholder_text(slide, 'title', title)
        
        chart_top = 1.0
        if subtitle:
            self._add_text_box(slide, subtitle, 0.4, 0.9, font_size=14, bold=True,
                             color=DataPackStyle.GREEN)
            chart_top = 1.4
            
        # Add chart image
        img_stream = BytesIO(chart_image)
        slide.shapes.add_picture(
            img_stream,
            Inches(0.4), Inches(chart_top),
            width=Inches(9.0)
        )
        
        if footnote:
            self._add_text_box(slide, footnote, 0.3, 6.5, font_size=6,
                             font_name=DataPackStyle.ANNOTATION_FONT,
                             color=DataPackStyle.ANNOTATION_GRAY)
        
        return slide
            
    def add_dual_chart_slide(self, title: str, 
                             chart1_image: bytes, chart1_title: str,
                             chart2_image: bytes, chart2_title: str):
        """Add slide with two charts stacked using two content layout"""
        # Try two content layout, fall back to title only
        layout = self.layout_two_content or self.layout_title_only
        slide = self.prs.slides.add_slide(layout)
        
        # Set title using placeholder
        self._set_placeholder_text(slide, 'title', title)
        
        # Top chart
        self._add_text_box(slide, chart1_title, 0.4, 0.9, font_size=12, bold=True,
                         color=DataPackStyle.GREEN)
        img1 = BytesIO(chart1_image)
        slide.shapes.add_picture(img1, Inches(0.4), Inches(1.2), width=Inches(9.0), height=Inches(2.5))
        
        # Bottom chart
        self._add_text_box(slide, chart2_title, 0.4, 3.9, font_size=12, bold=True,
                         color=DataPackStyle.GREEN)
        img2 = BytesIO(chart2_image)
        slide.shapes.add_picture(img2, Inches(0.4), Inches(4.2), width=Inches(9.0), height=Inches(2.5))
        
        return slide
        
    def add_top_customers_slide(self, title: str, data: pd.DataFrame):
        """Add top customers table slide using title only layout"""
        slide = self.prs.slides.add_slide(self.layout_title_only)
        
        # Set title using placeholder
        self._set_placeholder_text(slide, 'title', title)
        
        # Same table logic as P&L
        rows = min(len(data) + 1, 25)
        cols = min(len(data.columns), 8)
        
        table = slide.shapes.add_table(
            rows, cols,
            Inches(0.3), Inches(1.2),
            Inches(9.4), Inches(5.5)
        ).table
        
        # Header
        for j, col in enumerate(data.columns[:cols]):
            cell = table.cell(0, j)
            cell.text = str(col)
            cell.fill.solid()
            cell.fill.fore_color.rgb = DataPackStyle.NAVY
            p = cell.text_frame.paragraphs[0]
            p.font.color.rgb = DataPackStyle.WHITE
            p.font.size = Pt(8)
            p.font.bold = True
            p.font.name = DataPackStyle.BODY_FONT
            
        # Data
        for i, (idx, row) in enumerate(data.iterrows()):
            if i >= rows - 1:
                break
            for j, val in enumerate(row[:cols]):
                cell = table.cell(i + 1, j)
                cell.text = str(val) if pd.notna(val) else ""
                p = cell.text_frame.paragraphs[0]
                p.font.size = Pt(7)
                p.font.name = DataPackStyle.BODY_FONT
        
        return slide
                
    def save(self):
        """Save the presentation"""
        self.prs.save(self.output_path)
        return self.output_path


class ChartGenerator:
    """Generate charts matching Silver Oak style"""
    
    @staticmethod
    def setup_style():
        """Configure matplotlib for Silver Oak style"""
        plt.rcParams['font.family'] = 'sans-serif'
        plt.rcParams['font.sans-serif'] = ['Arial']
        plt.rcParams['axes.facecolor'] = 'white'
        plt.rcParams['figure.facecolor'] = 'white'
        plt.rcParams['axes.edgecolor'] = DataPackStyle.MPL_DARK
        plt.rcParams['axes.labelcolor'] = DataPackStyle.MPL_DARK
        plt.rcParams['xtick.color'] = DataPackStyle.MPL_DARK
        plt.rcParams['ytick.color'] = DataPackStyle.MPL_DARK
        
    @staticmethod
    def monthly_revenue_chart(dates: List, values: List, title: str = "") -> bytes:
        """Generate monthly revenue bar chart"""
        ChartGenerator.setup_style()
        
        fig, ax = plt.subplots(figsize=(9, 2.5), dpi=150)
        
        bars = ax.bar(range(len(values)), values, color=DataPackStyle.MPL_NAVY)
        
        ax.set_xticks(range(len(dates)))
        ax.set_xticklabels(dates, rotation=45, ha='right', fontsize=7)
        ax.set_ylabel('Revenue ($000s)', fontsize=9)
        ax.yaxis.set_major_formatter(plt.FuncFormatter(lambda x, p: f'${x:,.0f}'))
        
        ax.spines['top'].set_visible(False)
        ax.spines['right'].set_visible(False)
        
        plt.tight_layout()
        
        buf = BytesIO()
        plt.savefig(buf, format='png', bbox_inches='tight', dpi=150)
        plt.close()
        buf.seek(0)
        return buf.read()
        
    @staticmethod
    def ttm_revenue_chart(dates: List, values: List, title: str = "") -> bytes:
        """Generate TTM revenue line chart"""
        ChartGenerator.setup_style()
        
        fig, ax = plt.subplots(figsize=(9, 2.5), dpi=150)
        
        ax.plot(range(len(values)), values, color=DataPackStyle.MPL_GREEN, 
                linewidth=2, marker='o', markersize=4)
        ax.fill_between(range(len(values)), values, alpha=0.2, color=DataPackStyle.MPL_GREEN)
        
        ax.set_xticks(range(len(dates)))
        ax.set_xticklabels(dates, rotation=45, ha='right', fontsize=7)
        ax.set_ylabel('TTM Revenue ($M)', fontsize=9)
        ax.yaxis.set_major_formatter(plt.FuncFormatter(lambda x, p: f'${x:,.1f}'))
        
        ax.spines['top'].set_visible(False)
        ax.spines['right'].set_visible(False)
        
        plt.tight_layout()
        
        buf = BytesIO()
        plt.savefig(buf, format='png', bbox_inches='tight', dpi=150)
        plt.close()
        buf.seek(0)
        return buf.read()
        
    @staticmethod  
    def segment_breakdown_chart(segments: Dict[str, float]) -> bytes:
        """Generate segment breakdown pie/bar chart"""
        ChartGenerator.setup_style()
        
        fig, ax = plt.subplots(figsize=(4.5, 3), dpi=150)
        
        colors = [DataPackStyle.MPL_NAVY, DataPackStyle.MPL_GREEN, 
                  DataPackStyle.MPL_DARK, DataPackStyle.MPL_LIGHT]
        
        bars = ax.barh(list(segments.keys()), list(segments.values()), 
                       color=colors[:len(segments)])
        
        ax.set_xlabel('Revenue ($000s)', fontsize=9)
        ax.xaxis.set_major_formatter(plt.FuncFormatter(lambda x, p: f'${x:,.0f}'))
        
        ax.spines['top'].set_visible(False)
        ax.spines['right'].set_visible(False)
        
        plt.tight_layout()
        
        buf = BytesIO()
        plt.savefig(buf, format='png', bbox_inches='tight', dpi=150)
        plt.close()
        buf.seek(0)
        return buf.read()


class DataPackExcelGenerator:
    """Generate Data Pack backup Excel files"""
    
    def __init__(self, output_path: str):
        self.output_path = Path(output_path)
        self.writer = None
        self.sheets = {}
        
    def __enter__(self):
        self.writer = pd.ExcelWriter(self.output_path, engine='openpyxl')
        return self
        
    def __exit__(self, *args):
        if self.writer:
            self.writer.close()
            
    def add_sheet(self, name: str, df: pd.DataFrame, index: bool = True):
        """Add a sheet with data"""
        df.to_excel(self.writer, sheet_name=name[:31], index=index)
        
    def add_pl_sheet(self, name: str, data: pd.DataFrame):
        """Add formatted P&L sheet"""
        self.add_sheet(name, data)
        
    def add_customer_analysis(self, name: str, data: pd.DataFrame):
        """Add customer analysis sheet"""
        self.add_sheet(name, data, index=False)


def generate_datapack(
    company_name: str,
    financial_data: Dict[str, pd.DataFrame],
    customer_data: Dict[str, pd.DataFrame],
    output_dir: Path,
    date_str: str = None
) -> Dict[str, Path]:
    """
    Main function to generate a complete data pack
    
    Args:
        company_name: Name of the company
        financial_data: Dict of DataFrames with financial sheets
        customer_data: Dict of DataFrames with customer analysis
        output_dir: Directory for output files
        date_str: Date string for the pack (e.g., "October 2025")
        
    Returns:
        Dict with paths to generated files
    """
    output_dir = Path(output_dir)
    output_dir.mkdir(parents=True, exist_ok=True)
    
    timestamp = datetime.now().strftime("%Y%m%d_%H%M%S")
    
    # Generate PPT
    ppt_path = output_dir / f"{company_name.replace(' ', '_')}_Data_Pack_{timestamp}.pptx"
    ppt = DataPackPPTGenerator(ppt_path, company_name, date_str)
    
    # Title slide
    ppt.add_title_slide()
    
    # Agenda
    sections = [
        "Financial Trends",
        "By Segment", 
        "EBITDA Schedule",
        "Customer Trends"
    ]
    ppt.add_agenda_slide(sections, "Financial Trends")
    
    # P&L Summary slides
    if 'consolidated_pl' in financial_data:
        ppt.add_pl_summary_slide(f"Summary P&L – {company_name}", financial_data['consolidated_pl'])
        
    # Revenue charts
    if 'monthly_revenue' in financial_data:
        rev_data = financial_data['monthly_revenue']
        
        monthly_chart = ChartGenerator.monthly_revenue_chart(
            rev_data['date'].tolist(),
            rev_data['revenue'].tolist()
        )
        
        ttm_chart = ChartGenerator.ttm_revenue_chart(
            rev_data['date'].tolist(),
            rev_data['ttm_revenue'].tolist() if 'ttm_revenue' in rev_data else rev_data['revenue'].tolist()
        )
        
        ppt.add_dual_chart_slide(
            f"Monthly and Rolling TTM Revenue – {company_name}",
            monthly_chart, "Monthly Revenue",
            ttm_chart, "Rolling TTM Revenue"
        )
        
    # Segment analysis
    ppt.add_agenda_slide(sections, "By Segment")
    
    # Customer analysis
    ppt.add_agenda_slide(sections, "Customer Trends")
    
    if 'top_customers' in customer_data:
        ppt.add_top_customers_slide(
            f"Top Customers – {company_name}",
            customer_data['top_customers']
        )
    
    ppt.save()
    
    # Generate Excel backup
    excel_path = output_dir / f"{company_name.replace(' ', '_')}_Data_Pack_Backup_{timestamp}.xlsx"
    
    with DataPackExcelGenerator(excel_path) as excel:
        for sheet_name, df in financial_data.items():
            excel.add_sheet(sheet_name, df)
            
    # Generate Customer backup
    customer_excel_path = output_dir / f"{company_name.replace(' ', '_')}_Customer_Backup_{timestamp}.xlsx"
    
    with DataPackExcelGenerator(customer_excel_path) as excel:
        for sheet_name, df in customer_data.items():
            excel.add_sheet(sheet_name, df, index=False)
    
    return {
        'ppt': ppt_path,
        'data_backup': excel_path,
        'customer_backup': customer_excel_path
    }
